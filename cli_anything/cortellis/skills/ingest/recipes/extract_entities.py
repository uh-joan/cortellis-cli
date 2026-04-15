#!/usr/bin/env python3
"""
extract_entities.py — Find wiki entity mentions in free text.

Scans input text for known drug, target, company, and indication names
from the wiki INDEX.md. Returns matched entities as JSON.

Usage:
    python3 extract_entities.py <text_file>       # from file
    python3 extract_entities.py -                 # from stdin
    cat memo.md | python3 extract_entities.py -

Output (stdout): JSON list of {name, slug, type, count} sorted by count desc.
"""

import json
import os
import re
import sys


_TYPE_TO_DIR = {
    "indication": "indications",
    "company": "companies",
    "drug": "drugs",
    "target": "targets",
    "conference": "conferences",
}


def _read_aliases_from_frontmatter(path: str) -> list[str]:
    """Read the aliases: list from a wiki article's YAML frontmatter.

    Returns a list of alias strings, or [] if none found.
    Handles both list items (- alias) and inline scalars.
    """
    aliases = []
    in_frontmatter = False
    in_aliases = False

    with open(path, encoding="utf-8") as f:
        for i, raw_line in enumerate(f):
            line = raw_line.rstrip()
            if i == 0:
                in_frontmatter = line == "---"
                continue
            if in_frontmatter and line == "---":
                break
            if not in_frontmatter:
                break

            if line == "aliases:":
                in_aliases = True
                continue

            if in_aliases:
                if line.startswith("- "):
                    alias = line[2:].strip().strip("'\"")
                    if alias:
                        aliases.append(alias)
                elif line and not line.startswith(" ") and not line.startswith("-"):
                    in_aliases = False  # New YAML key — stop collecting

    return aliases


def load_wiki_entities(base_dir: str) -> list:
    """Parse wiki/INDEX.md and return list of {name, slug, type} dicts.

    Reads Markdown table rows of the form:
        | [Display Name](type/slug.md) | ... |
    and extracts name, slug, and article type.

    Also reads aliases: from each article's YAML frontmatter so that brand
    names, synonyms, and legal variants resolve to the canonical slug.
    E.g. "Wegovy" and "Ozempic" both resolve to slug "semaglutide".
    """
    index_path = os.path.join(base_dir, "wiki", "INDEX.md")
    if not os.path.exists(index_path):
        return []

    entities = []
    seen_slugs = set()
    current_type = None
    type_map = {
        "## Indications": "indication",
        "## Companies": "company",
        "## Drugs": "drug",
        "## Targets": "target",
        "## Conferences": "conference",
        "## Internal": "internal",
    }

    with open(index_path, encoding="utf-8") as f:
        for line in f:
            line = line.rstrip()
            # Detect section headers
            for header, etype in type_map.items():
                if line.startswith(header):
                    current_type = etype
                    break

            if not current_type:
                continue

            # Match table rows with wikilinks: | [Name](path/slug.md) |
            m = re.search(r'\[([^\]]+)\]\(([^)]+\.md)\)', line)
            if not m:
                continue
            name = m.group(1).strip()
            path = m.group(2).strip()
            # Extract slug from path: "drugs/semaglutide.md" → "semaglutide"
            slug = os.path.splitext(os.path.basename(path))[0]

            if slug in seen_slugs or not name:
                continue
            seen_slugs.add(slug)
            entities.append({"name": name, "slug": slug, "type": current_type})

    # Extend with aliases from article frontmatter
    alias_entries = []
    seen_alias_keys: set[tuple] = set()
    wiki_dir = os.path.join(base_dir, "wiki")

    for entity in entities:
        etype = entity["type"]
        type_dir = _TYPE_TO_DIR.get(etype)
        if not type_dir:
            continue
        art_path = os.path.join(wiki_dir, type_dir, f"{entity['slug']}.md")
        if not os.path.exists(art_path):
            continue
        try:
            for alias in _read_aliases_from_frontmatter(art_path):
                key = (alias.lower(), entity["slug"])
                if key in seen_alias_keys:
                    continue
                seen_alias_keys.add(key)
                alias_entries.append({
                    "name": alias,
                    "slug": entity["slug"],
                    "type": entity["type"],
                })
        except Exception:
            continue

    # Also load from wiki/entity_aliases.json (abbreviations, short forms)
    alias_entries.extend(_load_json_aliases(wiki_dir, entities, seen_alias_keys))

    return entities + alias_entries


def _load_json_aliases(wiki_dir: str, entities: list, seen: set) -> list:
    """Load alias→slug mappings from wiki/entity_aliases.json.

    File format: {"t2d": "diabetes-mellitus", "mash": "metabolic-dysfunction-..."}
    Keys are lowercase; values are entity slugs that must exist in INDEX.md.
    """
    alias_path = os.path.join(wiki_dir, "entity_aliases.json")
    if not os.path.exists(alias_path):
        return []

    slug_to_entity = {e["slug"]: e for e in entities}
    extras = []
    try:
        with open(alias_path, encoding="utf-8") as f:
            alias_map = json.load(f)
        for alias, slug in alias_map.items():
            entity = slug_to_entity.get(slug)
            if not entity:
                continue
            key = (alias.lower(), slug)
            if key in seen:
                continue
            seen.add(key)
            extras.append({
                "name": alias,
                "slug": slug,
                "type": entity["type"],
            })
    except Exception:
        pass
    return extras


def scan_text(text: str, entities: list) -> list:
    """Find entity name mentions in text. Returns list of matched entities with count.

    Matching is case-insensitive and whole-word. Shorter names that are substrings
    of longer matched names are suppressed to avoid double-counting.
    """
    text.lower()
    matches = []

    # Sort by name length descending so longer names match first
    sorted_entities = sorted(entities, key=lambda e: -len(e["name"]))
    matched_spans = []  # track (start, end) of consumed spans

    for entity in sorted_entities:
        name = entity["name"]
        pattern = r'\b' + re.escape(name) + r'\b'
        found = list(re.finditer(pattern, text, re.IGNORECASE))
        if not found:
            continue

        # Check if all matches are within already-matched spans
        new_matches = []
        for m in found:
            span = (m.start(), m.end())
            overlaps = any(
                not (span[1] <= s[0] or span[0] >= s[1])
                for s in matched_spans
            )
            if not overlaps:
                new_matches.append(span)

        if new_matches:
            matched_spans.extend(new_matches)
            matches.append({
                "name": name,
                "slug": entity["slug"],
                "type": entity["type"],
                "count": len(new_matches),
            })

    # Sort by count descending, then name ascending
    matches.sort(key=lambda x: (-x["count"], x["name"]))
    return matches


def inject_wikilinks(text: str, matched_entities: list) -> str:
    """Replace entity name mentions in body text with [[slug|Name]] wikilinks.

    Processes entities from longest name to shortest to avoid partial substitutions.
    Only replaces the first occurrence to avoid over-linking.
    """
    sorted_matches = sorted(matched_entities, key=lambda e: -len(e["name"]))
    for entity in sorted_matches:
        name = entity["name"]
        slug = entity["slug"]
        wikilink = f"[[{slug}|{name}]]"
        # Replace all occurrences (case-insensitive)
        pattern = r'\b' + re.escape(name) + r'\b'
        text = re.sub(pattern, wikilink, text, flags=re.IGNORECASE)
    return text


def extract_text_from_file(path: str) -> str:
    """Extract plain text from a file.

    Supported formats:
      .md / .txt / .markdown / .rst  — direct read
      .csv                           — direct read (UTF-8-sig, handles Excel BOM)
      .pdf                           — pdftotext (requires poppler)
      .pptx                          — python-pptx
      .xlsx / .xlsm                  — openpyxl (text cells only)
      .xlsb                          — not supported (binary Excel)
    """
    ext = os.path.splitext(path)[1].lower()

    if ext in (".md", ".txt", ".markdown", ".rst", ".csv"):
        # utf-8-sig strips BOM present in Excel-exported CSVs
        with open(path, encoding="utf-8-sig") as f:
            return f.read()

    if ext == ".pdf":
        import subprocess
        result = subprocess.run(
            ["pdftotext", "-layout", path, "-"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            raise ValueError(f"pdftotext failed: {result.stderr.strip()}")
        return result.stdout

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            chunks.append(text)
        return "\n".join(chunks)

    if ext in (".xlsx", ".xlsm"):
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        chunks = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell and isinstance(cell, str) and cell.strip():
                        chunks.append(cell.strip())
        return "\n".join(chunks)

    if ext == ".xlsb":
        raise ValueError("Binary Excel (.xlsb) is not supported. Convert to .xlsx first.")

    raise ValueError(f"Unsupported file type: {ext}. Supported: .md, .txt, .pdf, .pptx, .xlsx, .xlsm")


def cortellis_ner_entities(text: str, base_dir: str) -> list:
    """Use Cortellis NER API to extract entities, then resolve against wiki slugs.

    Falls back silently to [] on auth failure or network error.
    Returns same shape as scan_text(): [{name, slug, type, count}]
    """
    try:
        from dotenv import load_dotenv
        load_dotenv()
        from cli_anything.cortellis.core.client import CortellisClient
        from cli_anything.cortellis.core import ner as _ner

        client = CortellisClient()
        result = _ner.match(client, text)
        raw_entities = (
            result.get("NamedEntityRecognition", {})
                  .get("Entities", {})
                  .get("Entity", [])
        )
        if isinstance(raw_entities, dict):
            raw_entities = [raw_entities]
        if not raw_entities:
            return []

        # Build wiki slug lookup from INDEX
        wiki_entities = load_wiki_entities(base_dir)
        name_to_entity = {e["name"].lower(): e for e in wiki_entities}
        slug_to_entity = {e["slug"]: e for e in wiki_entities}

        _TYPE_MAP = {"Drug": "drug", "Target": "target", "Company": "company",
                     "Indication": "indication", "Disease": "indication"}

        seen = {}
        for ent in raw_entities:
            synonym = ent.get("@synonym", "").strip()
            ner_type = _TYPE_MAP.get(ent.get("@type", ""), "")
            if not synonym or not ner_type:
                continue

            # Resolve to wiki slug: exact name match first, then slugify
            wiki_ent = name_to_entity.get(synonym.lower())
            if not wiki_ent:
                from cli_anything.cortellis.utils.wiki import slugify
                candidate_slug = slugify(synonym)
                wiki_ent = slug_to_entity.get(candidate_slug)
            if not wiki_ent:
                continue  # Entity not in wiki yet — skip

            key = wiki_ent["slug"]
            if key in seen:
                seen[key]["count"] += 1
            else:
                seen[key] = {
                    "name": wiki_ent["name"],
                    "slug": wiki_ent["slug"],
                    "type": wiki_ent["type"],
                    "count": 1,
                }

        matched = sorted(seen.values(), key=lambda x: (-x["count"], x["name"]))
        print(f"[ner] Cortellis NER: {len(matched)} wiki entities resolved", file=sys.stderr)
        return matched

    except Exception as exc:
        print(f"[ner] Cortellis NER unavailable ({exc.__class__.__name__}), using wiki-index fallback", file=sys.stderr)
        return []


def main():
    if len(sys.argv) < 2:
        print("Usage: extract_entities.py <text_file|->\n", file=sys.stderr)
        print("  <text_file>  Path to .md/.txt file", file=sys.stderr)
        print("  -            Read from stdin", file=sys.stderr)
        sys.exit(1)

    source = sys.argv[1]
    if source == "-":
        text = sys.stdin.read()
    else:
        if not os.path.exists(source):
            print(f"Error: file not found: {source}", file=sys.stderr)
            sys.exit(1)
        text = extract_text_from_file(source)

    if not text.strip():
        print("[]")
        return

    base_dir = os.getcwd()

    # Try Cortellis NER first; fall back to wiki-index scan
    matched = cortellis_ner_entities(text, base_dir)
    if not matched:
        wiki_entities = load_wiki_entities(base_dir)
        if not wiki_entities:
            print(f"[warn] No entities in {base_dir}/wiki/INDEX.md", file=sys.stderr)
            print("[]")
            return
        matched = scan_text(text, wiki_entities)

    print(json.dumps(matched, indent=2))

    if matched:
        print(f"Found {len(matched)} entities:", file=sys.stderr)
        for e in matched[:10]:
            print(f"  [{e['type']}] {e['name']} ({e['count']}x)", file=sys.stderr)
        if len(matched) > 10:
            print(f"  ... and {len(matched) - 10} more", file=sys.stderr)
    else:
        print("No known wiki entities found in text.", file=sys.stderr)


if __name__ == "__main__":
    main()
