#!/usr/bin/env python3
"""
export_pptx.py — Export landscape analysis as PowerPoint deck.

Reads scored CSVs and markdown reports from a landscape directory
and produces a professional 16:9 slide deck.

Usage: python3 export_pptx.py <landscape_dir> [indication_name] [--output PATH]
"""

import os
import re
import sys
from datetime import datetime, timezone

# Allow running as standalone script
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "..", "..", "..", ".."))

from cli_anything.cortellis.utils.data_helpers import (
    read_csv_safe,
    safe_float,
    safe_int,
    read_json_safe,
    read_md_safe,
    count_csv_rows,
)
from cli_anything.cortellis.utils.wiki import slugify

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x1B, 0x3A, 0x5C)
PHARMA_BLUE = RGBColor(0x00, 0x7B, 0xC0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)



# ---------------------------------------------------------------------------
# Data loading (same as compile_dossier)
# ---------------------------------------------------------------------------

PHASE_FILES = ["launched", "phase3", "phase2", "phase1", "discovery", "other"]
PHASE_LABELS = {
    "launched": "Launched",
    "phase3": "Phase 3",
    "phase2": "Phase 2",
    "phase1": "Phase 1",
    "discovery": "Discovery",
    "other": "Other",
}


def load_phase_counts(landscape_dir):
    counts = {}
    total = 0
    for phase in PHASE_FILES:
        n = count_csv_rows(landscape_dir, f"{phase}.csv")
        counts[phase] = n
        total += n
    counts["total"] = total
    return counts


def load_strategic_scores(landscape_dir):
    return read_csv_safe(os.path.join(landscape_dir, "strategic_scores.csv"))


def load_mechanism_scores(landscape_dir):
    return read_csv_safe(os.path.join(landscape_dir, "mechanism_scores.csv"))


def load_opportunity_matrix(landscape_dir):
    return read_csv_safe(os.path.join(landscape_dir, "opportunity_matrix.csv"))


def load_freshness(landscape_dir):
    return read_json_safe(os.path.join(landscape_dir, "freshness.json"))


def detect_preset(landscape_dir):
    trail = read_json_safe(os.path.join(landscape_dir, "audit_trail.json"))
    if isinstance(trail, list) and trail:
        for entry in reversed(trail):
            preset = entry.get("preset", {})
            if preset and preset.get("name"):
                return preset["name"]
    return "default"


# ---------------------------------------------------------------------------
# Presentation factory
# ---------------------------------------------------------------------------

def create_presentation() -> Presentation:
    """Create a blank 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------

def _blank_layout(prs):
    """Return a blank slide layout (last one is usually blank)."""
    return prs.slide_layouts[6]  # index 6 = blank in default theme


def _add_footer(slide, date_str):
    """Add a small footer text box at the bottom of every slide."""
    left = Inches(0.4)
    top = Inches(7.1)
    width = Inches(12.5)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"Source: Cortellis  |  {date_str}"
    run.font.size = Pt(8)
    run.font.color.rgb = MID_GRAY
    run.font.name = "Calibri"


def _add_slide_title(slide, title_text, subtitle_text=None):
    """Add a title bar at the top of the slide."""
    # Navy bar
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(1.1)
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1; pptx constant for rect
        left, top, width, height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.size = Pt(13)
        r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        r2.font.name = "Calibri"


def _add_body_textbox(slide, text_lines, left, top, width, height, font_size=12):
    """Add a multi-line body textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in text_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Calibri"


def _style_table_header_row(table):
    """Style the first row of a table as a navy header."""
    from pptx.oxml.ns import qn
    from lxml import etree
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.name = "Calibri"
                run.font.size = Pt(10)


def _set_cell_text(cell, text, font_size=9, bold=False, color=None):
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # Clear existing runs
    for run in p.runs:
        run.text = ""
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Calibri"
    if color:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = DARK_GRAY


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(_blank_layout(prs))

    # Full navy background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Pharma blue accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(3.6), Inches(13.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PHARMA_BLUE
    accent.line.fill.background()

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(1.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.3), Inches(11), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    r2.font.name = "Calibri"

    # Date
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(11), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = date_str
    r3.font.size = Pt(14)
    r3.font.color.rgb = MID_GRAY
    r3.font.name = "Calibri"

    # Footer
    _add_footer(slide, date_str)
    return slide


def add_summary_slide(prs, phases, scores, deal_count):
    """Slide 2: Executive Summary."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    _add_slide_title(slide, "Executive Summary")
    _add_footer(slide, date_str)

    # Key metrics boxes — 4 across top
    metrics = [
        ("Total Drugs", str(phases.get("total", 0))),
        ("Deals", str(deal_count)),
        ("Companies\n(top tier)", str(len([r for r in scores if r.get("cpi_tier") == "A"]))),
        ("Launched", str(phases.get("launched", 0))),
    ]
    box_w = Inches(2.8)
    box_h = Inches(1.3)
    for i, (label, value) in enumerate(metrics):
        lft = Inches(0.4 + i * 3.1)
        top = Inches(1.3)
        # Background box
        box = slide.shapes.add_shape(1, lft, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = PHARMA_BLUE

        # Value
        vbox = slide.shapes.add_textbox(lft + Inches(0.1), top + Inches(0.05), box_w - Inches(0.2), Inches(0.75))
        vp = vbox.text_frame.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vrun = vp.add_run()
        vrun.text = value
        vrun.font.size = Pt(28)
        vrun.font.bold = True
        vrun.font.color.rgb = NAVY
        vrun.font.name = "Calibri"

        # Label
        lbox = slide.shapes.add_textbox(lft + Inches(0.1), top + Inches(0.75), box_w - Inches(0.2), Inches(0.45))
        lp = lbox.text_frame.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lrun = lp.add_run()
        lrun.text = label
        lrun.font.size = Pt(10)
        lrun.font.color.rgb = DARK_GRAY
        lrun.font.name = "Calibri"

    # Phase breakdown bullets
    bullet_lines = []
    for phase_key in ["launched", "phase3", "phase2", "phase1", "discovery"]:
        label = PHASE_LABELS[phase_key]
        bullet_lines.append(f"  {label}: {phases.get(phase_key, 0)} drugs")

    # Top company
    if scores:
        top = scores[0]
        bullet_lines.insert(0, f"  Top company: {top.get('company', '-')} (CPI {safe_float(top.get('cpi_score')):.1f})")

    _add_body_textbox(
        slide, bullet_lines,
        Inches(0.5), Inches(2.8),
        Inches(12.3), Inches(4.0),
        font_size=13,
    )
    return slide


def add_pipeline_chart_slide(prs, phases):
    """Slide 3: Pipeline Overview bar chart."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    _add_slide_title(slide, "Pipeline Overview", "Drug count by development phase")
    _add_footer(slide, date_str)

    chart_data = CategoryChartData()
    display_phases = ["launched", "phase3", "phase2", "phase1", "discovery"]
    categories = [PHASE_LABELS[p] for p in display_phases]
    values = [phases.get(p, 0) for p in display_phases]

    chart_data.categories = categories
    chart_data.add_series("Drugs", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(1.0), Inches(1.3),
        Inches(11.3), Inches(5.8),
        chart_data,
    ).chart

    chart.has_legend = False
    chart.has_title = False

    # Style the series fill
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = PHARMA_BLUE

    # Add data labels
    series.data_labels.show_value = True
    series.data_labels.font.size = Pt(11)
    series.data_labels.font.name = "Calibri"

    return slide


def add_competitive_table_slide(prs, scores, max_rows=15):
    """Slide 4: Competitive Landscape table."""
    if not scores:
        return None

    slide = prs.slides.add_slide(_blank_layout(prs))
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    _add_slide_title(slide, "Competitive Landscape", f"Top {min(max_rows, len(scores))} companies by CPI score")
    _add_footer(slide, date_str)

    rows_data = scores[:max_rows]
    n_rows = len(rows_data) + 1  # +1 for header
    cols = 5
    headers = ["Rank", "Company", "Tier", "CPI Score", "Position"]
    col_widths = [Inches(0.6), Inches(4.5), Inches(0.8), Inches(1.2), Inches(1.8)]

    left = Inches(0.4)
    top = Inches(1.25)
    width = sum(col_widths)
    row_h = Inches(0.28)
    height = row_h * n_rows

    table = slide.shapes.add_table(n_rows, cols, left, top, width, height).table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # Header row
    for ci, h in enumerate(headers):
        _set_cell_text(table.cell(0, ci), h, font_size=10, bold=True, color=WHITE)
    _style_table_header_row(table)

    # Data rows
    for ri, row in enumerate(rows_data, 1):
        cpi = safe_float(row.get("cpi_score"))
        tier = row.get("cpi_tier", "-")
        _set_cell_text(table.cell(ri, 0), str(ri), font_size=9)
        _set_cell_text(table.cell(ri, 1), row.get("company", "-"), font_size=9)
        _set_cell_text(table.cell(ri, 2), tier, font_size=9)
        _set_cell_text(table.cell(ri, 3), f"{cpi:.1f}", font_size=9)
        _set_cell_text(table.cell(ri, 4), row.get("position", "-"), font_size=9)

        # Alternating row background
        if ri % 2 == 0:
            for ci in range(cols):
                table.cell(ri, ci).fill.solid()
                table.cell(ri, ci).fill.fore_color.rgb = LIGHT_GRAY

    return slide


def add_mechanism_slide(prs, mechanisms, max_rows=10):
    """Slide 5: Mechanism Analysis table."""
    if not mechanisms:
        return None

    slide = prs.slides.add_slide(_blank_layout(prs))
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    _add_slide_title(slide, "Mechanism Analysis", f"Top {min(max_rows, len(mechanisms))} mechanisms by crowding index")
    _add_footer(slide, date_str)

    rows_data = mechanisms[:max_rows]
    n_rows = len(rows_data) + 1
    cols = 6
    headers = ["Mechanism", "Active", "Launched", "Phase 3", "Companies", "Crowding Index"]
    col_widths = [Inches(4.0), Inches(0.85), Inches(0.95), Inches(0.85), Inches(1.2), Inches(1.6)]

    left = Inches(0.4)
    top = Inches(1.25)
    width = sum(col_widths)
    row_h = Inches(0.32)
    height = row_h * n_rows

    table = slide.shapes.add_table(n_rows, cols, left, top, width, height).table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ci, h in enumerate(headers):
        _set_cell_text(table.cell(0, ci), h, font_size=10, bold=True, color=WHITE)
    _style_table_header_row(table)

    for ri, row in enumerate(rows_data, 1):
        _set_cell_text(table.cell(ri, 0), row.get("mechanism", "-"), font_size=9)
        _set_cell_text(table.cell(ri, 1), str(safe_int(row.get("active_count"))), font_size=9)
        _set_cell_text(table.cell(ri, 2), str(safe_int(row.get("launched"))), font_size=9)
        _set_cell_text(table.cell(ri, 3), str(safe_int(row.get("phase3"))), font_size=9)
        _set_cell_text(table.cell(ri, 4), str(safe_int(row.get("company_count"))), font_size=9)
        _set_cell_text(table.cell(ri, 5), f"{safe_int(row.get('crowding_index')):,}", font_size=9)

        if ri % 2 == 0:
            for ci in range(cols):
                table.cell(ri, ci).fill.solid()
                table.cell(ri, ci).fill.fore_color.rgb = LIGHT_GRAY

    return slide


def add_deals_slide(prs, landscape_dir):
    """Slide 6: Deal Landscape."""
    deal_count = count_csv_rows(landscape_dir, "deals.csv")
    deals_md = read_md_safe(os.path.join(landscape_dir, "deals_analytics.md"))
    deal_comps_md = read_md_safe(os.path.join(landscape_dir, "deal_comps.md"))

    slide = prs.slides.add_slide(_blank_layout(prs))
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    _add_slide_title(slide, "Deal Landscape", f"{deal_count} deals in dataset")
    _add_footer(slide, date_str)

    # Parse deal type counts from deals.csv
    deals_rows = read_csv_safe(os.path.join(landscape_dir, "deals.csv"))
    type_counts = {}
    for d in deals_rows:
        dtype = d.get("deal_type") or d.get("type") or d.get("activity_type") or "Unknown"
        type_counts[dtype] = type_counts.get(dtype, 0) + 1
    top_types = sorted(type_counts.items(), key=lambda x: x[1], reverse=True)[:6]

    lines = []
    lines.append(f"Total deals: {deal_count}")
    lines.append("")
    if top_types:
        lines.append("Top deal types:")
        for dtype, cnt in top_types:
            lines.append(f"  {dtype}: {cnt}")

    # Extract a few lines from deals_analytics.md if available
    if deals_md:
        lines.append("")
        lines.append("Analytics summary:")
        for line in deals_md.splitlines()[:8]:
            stripped = line.strip()
            if stripped and not stripped.startswith("#"):
                lines.append(f"  {stripped}")

    if deal_comps_md:
        lines.append("")
        lines.append("Financial terms summary:")
        for line in deal_comps_md.splitlines()[:6]:
            stripped = line.strip()
            if stripped and not stripped.startswith("#"):
                lines.append(f"  {stripped}")

    _add_body_textbox(
        slide, lines,
        Inches(0.5), Inches(1.35),
        Inches(12.3), Inches(5.7),
        font_size=12,
    )
    return slide


def add_opportunity_slide(prs, opportunities):
    """Slide 7: Opportunity Assessment."""
    if not opportunities:
        return None

    slide = prs.slides.add_slide(_blank_layout(prs))
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    _add_slide_title(slide, "Opportunity Assessment", "White space vs crowded mechanisms")
    _add_footer(slide, date_str)

    white_space = [r for r in opportunities if r.get("status") in ("White Space", "Emerging")]
    crowded = [r for r in opportunities if r.get("status") == "Crowded Pipeline"]

    # Build two-section table: white space rows + crowded rows
    all_rows = []
    for r in white_space[:8]:
        all_rows.append(("white", r))
    for r in crowded[:7]:
        all_rows.append(("crowded", r))

    if not all_rows:
        return slide

    n_rows = len(all_rows) + 1
    cols = 5
    headers = ["Mechanism", "Status", "Total Drugs", "Companies", "Opportunity Score"]
    col_widths = [Inches(4.5), Inches(1.5), Inches(1.3), Inches(1.3), Inches(1.9)]

    left = Inches(0.4)
    top = Inches(1.25)
    width = sum(col_widths)
    row_h = Inches(0.3)
    height = row_h * n_rows

    table = slide.shapes.add_table(n_rows, cols, left, top, width, height).table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ci, h in enumerate(headers):
        _set_cell_text(table.cell(0, ci), h, font_size=10, bold=True, color=WHITE)
    _style_table_header_row(table)

    for ri, (row_type, row) in enumerate(all_rows, 1):
        status = row.get("status", "-")
        opp_score = safe_float(row.get("opportunity_score"))
        _set_cell_text(table.cell(ri, 0), row.get("mechanism", "-"), font_size=9)
        _set_cell_text(table.cell(ri, 1), status, font_size=9)
        _set_cell_text(table.cell(ri, 2), str(safe_int(row.get("total"))), font_size=9)
        _set_cell_text(table.cell(ri, 3), str(safe_int(row.get("companies"))), font_size=9)
        _set_cell_text(table.cell(ri, 4), f"{opp_score:.4f}", font_size=9)

        # Color code: green tint for white space, light for crowded
        if row_type == "white":
            for ci in range(cols):
                table.cell(ri, ci).fill.solid()
                table.cell(ri, ci).fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
        elif ri % 2 == 0:
            for ci in range(cols):
                table.cell(ri, ci).fill.solid()
                table.cell(ri, ci).fill.fore_color.rgb = LIGHT_GRAY

    return slide


def add_sources_slide(prs, landscape_dir, preset):
    """Slide 8: Data Sources."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    date_str = datetime.now(timezone.utc).strftime("%B %Y")
    _add_slide_title(slide, "Data Sources")
    _add_footer(slide, date_str)

    freshness = load_freshness(landscape_dir)
    now_str = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

    lines = [
        f"Source directory:   {os.path.abspath(landscape_dir)}",
        f"Freshness level:    {freshness.get('staleness_level', 'unknown')}",
        f"Data computed:      {freshness.get('computed_at_utc', 'unknown')}",
        f"Preset:             {preset}",
        f"Report compiled:    {now_str}",
        "",
        "Data files used:",
    ]
    for phase in PHASE_FILES:
        path = os.path.join(landscape_dir, f"{phase}.csv")
        if os.path.exists(path):
            n = count_csv_rows(landscape_dir, f"{phase}.csv")
            lines.append(f"  {phase}.csv  ({n} drugs)")

    for fname in ["deals.csv", "strategic_scores.csv", "mechanism_scores.csv", "opportunity_matrix.csv"]:
        path = os.path.join(landscape_dir, fname)
        if os.path.exists(path):
            n = count_csv_rows(landscape_dir, fname)
            lines.append(f"  {fname}  ({n} rows)")

    _add_body_textbox(
        slide, lines,
        Inches(0.5), Inches(1.35),
        Inches(12.3), Inches(5.7),
        font_size=12,
    )
    return slide


# ---------------------------------------------------------------------------
# Main export
# ---------------------------------------------------------------------------

def export_pptx(landscape_dir, indication_name, output_path=None) -> str:
    """Build and save the PowerPoint deck. Returns the output path."""
    if not os.path.isdir(landscape_dir):
        raise FileNotFoundError(f"Landscape directory not found: {landscape_dir}")

    if not indication_name:
        indication_name = os.path.basename(landscape_dir).replace("-", " ").title()

    # Derive output path
    if not output_path:
        slug = slugify(indication_name)
        output_path = os.path.join(landscape_dir, f"{slug}-landscape.pptx")

    # Load data
    phases = load_phase_counts(landscape_dir)
    scores = load_strategic_scores(landscape_dir)
    mechanisms = load_mechanism_scores(landscape_dir)
    opportunities = load_opportunity_matrix(landscape_dir)
    preset = detect_preset(landscape_dir)
    deal_count = count_csv_rows(landscape_dir, "deals.csv")

    prs = create_presentation()
    slide_count = 0

    # Slide 1: Title
    add_title_slide(prs, indication_name, "Competitive Landscape Report")
    slide_count += 1

    # Slide 2: Executive Summary
    add_summary_slide(prs, phases, scores, deal_count)
    slide_count += 1

    # Slide 3: Pipeline chart (always, even if all zeros)
    add_pipeline_chart_slide(prs, phases)
    slide_count += 1

    # Slide 4: Competitive table (skip if no data)
    if scores:
        add_competitive_table_slide(prs, scores, max_rows=15)
        slide_count += 1

    # Slide 5: Mechanism analysis (skip if no data)
    if mechanisms:
        add_mechanism_slide(prs, mechanisms, max_rows=10)
        slide_count += 1

    # Slide 6: Deal landscape (always — graceful when no deals)
    add_deals_slide(prs, landscape_dir)
    slide_count += 1

    # Slide 7: Opportunity assessment (skip if no data)
    if opportunities:
        add_opportunity_slide(prs, opportunities)
        slide_count += 1

    # Slide 8: Data sources
    add_sources_slide(prs, landscape_dir, preset)
    slide_count += 1

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    if len(sys.argv) < 2:
        print(
            "Usage: export_pptx.py <landscape_dir> [indication_name] [--output PATH]",
            file=sys.stderr,
        )
        sys.exit(1)

    landscape_dir = sys.argv[1]
    indication_name = None
    output_path = None

    i = 2
    while i < len(sys.argv):
        arg = sys.argv[i]
        if arg == "--output" and i + 1 < len(sys.argv):
            output_path = sys.argv[i + 1]
            i += 2
        elif not arg.startswith("--"):
            indication_name = arg
            i += 1
        else:
            i += 1

    if not os.path.isdir(landscape_dir):
        print(f"Error: landscape directory not found: {landscape_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"Exporting {indication_name or landscape_dir} to PowerPoint...")
    out = export_pptx(landscape_dir, indication_name, output_path)
    print(f"Saved: {out}")

    # Verify slide count
    prs = Presentation(out)
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
